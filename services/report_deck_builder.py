"""Leadership deck (PPTX) export for the Report Summary Contract v1.

Renders a completed audit's `report_summary` into a short, boardroom-ready
deck: cover (score + band + as-of date), verdict + competitive snapshot, an
optional LLM-written executive summary, one slide per top action (with the
measured prompt evidence as captions), and a methodology/honesty slide.

Pricing (PR-4 decision): the ONLY LLM step is the executive summary; it bills
on ACTUAL token usage at DECK_TOKEN_PRICE_MULTIPLE (1.6x measured token COGS —
deliberately above the 1.2 probe flat_multiple: the deck is a premium,
leadership-facing artifact). Everything else is deterministic rendering of the
contract, so a deck without the LLM step (no key / LLM down) costs 0 credits.

Honesty rules carried over from the contract: every number and finding on the
deck is verbatim contract data; the LLM bullets may only rephrase for an
executive audience and are prompted to add no new facts, numbers, or claims;
the deck prints its "data as of" date and the run's honest_limits.

python-pptx is imported lazily (mirrors the weasyprint pattern in
audit_html_renderer): build_report_deck returns None when it isn't installed
so the route can 503 with a clear message instead of crashing at import time.
"""

from __future__ import annotations

import json
import logging
from decimal import Decimal
from typing import Any, Dict, List, Mapping, Optional, Tuple

import httpx

from config.settings import settings

logger = logging.getLogger(__name__)

# Customer price = measured token COGS x this multiple (user-set pricing:
# 1.6x actual consumption), converted to credits by credits_for_tokens.
DECK_TOKEN_PRICE_MULTIPLE = Decimal("1.6")

DECK_LLM_PROVIDER = "deepseek"
_EXEC_SUMMARY_MAX_TOKENS = 400
_EXEC_SUMMARY_TIMEOUT_S = 30.0

_EXEC_SUMMARY_SYSTEM = (
    "You write one executive-summary slide for a brand's AI-readiness audit "
    "deck. Audience: marketing leadership deciding what to fund next. Input "
    "is the audit's summary JSON. Write 3-4 short bullets (max ~18 words "
    "each): the current state, what it costs the brand, and what to do first. "
    "STRICT: use only facts, numbers, and names present in the JSON — never "
    "invent, extrapolate, or soften. Plain business English, no jargon, no "
    "emoji. Output one bullet per line, no numbering, no markdown."
)

# Midnight-executive palette (skill-vetted): navy dominates the cover +
# closing, ice-blue accents, white content slides. Safe fonts only.
_NAVY = (0x1E, 0x27, 0x61)
_ICE = (0xCA, 0xDC, 0xFC)
_WHITE = (0xFF, 0xFF, 0xFF)
_INK = (0x21, 0x21, 0x21)
_MUTED = (0x5A, 0x5F, 0x73)
_RED = (0xB4, 0x23, 0x18)
_AMBER = (0xA1, 0x62, 0x07)
_GREEN = (0x1E, 0x6B, 0x3C)

_BAND_DECK_LABEL = {
    "needs_work": ("Needs work", _AMBER),
    "pass": ("Pass", _NAVY),
    "good": ("Good", _GREEN),
    "excellent": ("Excellent", _GREEN),
}

_SEVERITY_COLOR = {"high": _RED, "medium": _AMBER, "info": _MUTED}


def _as_dict(value: Any) -> Dict[str, Any]:
    return value if isinstance(value, dict) else {}


def _as_list(value: Any) -> List[Any]:
    return value if isinstance(value, list) else []


async def generate_executive_summary(
    summary: Mapping[str, Any],
) -> Optional[Tuple[List[str], int, int]]:
    """LLM executive-summary bullets grounded ONLY in the contract JSON.

    Returns (bullets, prompt_tokens, completion_tokens), or None when no API
    key is configured / the provider returned nothing usable — the deck then
    ships without this slide and nothing is billed. Raises on transport errors
    (caller catches). Isolated so tests can monkeypatch it.
    """
    api_key = (settings.deepseek_api_key or "").strip()
    if not api_key:
        return None
    # Only the fields leadership copy may draw from — strips sku_summaries etc.
    # Share-of-voice + subscores are included so the summary can reference the
    # competitive rank and which dimension is the drag (the whole point of the
    # deck) — before this they were absent, so the LLM couldn't name either.
    _score = _as_dict(summary.get("score"))
    grounding = {
        "score": {
            "display": _score.get("display"),
            "scale_max": _score.get("scale_max"),
            "band": _score.get("band"),
            "subscores": _as_list(_score.get("subscores")),
            "weakest_dimension": _as_dict(_score.get("weakest_dimension")),
        },
        "verdict": _as_dict(summary.get("verdict")),
        "share_of_voice": _as_dict(summary.get("share_of_voice")),
        "since_last_audit": {
            k: v
            for k, v in _as_dict(summary.get("since_last_audit")).items()
            if k in ("headline", "movements")
        },
        "top_findings": _as_list(summary.get("top_findings")),
        "top_actions": [
            {
                k: v
                for k, v in _as_dict(a).items()
                if k in ("headline", "why_this_first", "first_move", "sku_title")
            }
            for a in _as_list(summary.get("top_actions"))
        ],
        "competitive_snapshot": _as_dict(summary.get("competitive_snapshot")),
    }
    payload = {
        "model": settings.deepseek_model,
        "messages": [
            {"role": "system", "content": _EXEC_SUMMARY_SYSTEM},
            {"role": "user", "content": json.dumps(grounding, ensure_ascii=False)},
        ],
        "temperature": 0.2,
        "max_tokens": _EXEC_SUMMARY_MAX_TOKENS,
        "stream": False,
    }
    headers = {
        "Authorization": f"Bearer {api_key}",
        "content-type": "application/json",
    }
    base_url = settings.deepseek_api_base_url.rstrip("/")
    async with httpx.AsyncClient(timeout=_EXEC_SUMMARY_TIMEOUT_S) as client:
        response = await client.post(
            f"{base_url}/v1/chat/completions", json=payload, headers=headers
        )
        response.raise_for_status()
        data = response.json()
    content = (data["choices"][0]["message"]["content"] or "").strip()
    bullets = [
        line.strip().lstrip("-•*").strip()
        for line in content.splitlines()
        if line.strip()
    ][:4]
    if not bullets:
        return None
    usage = data.get("usage") or {}
    return (
        bullets,
        int(usage.get("prompt_tokens") or 0),
        int(usage.get("completion_tokens") or 0),
    )


def _fmt_display(score: Mapping[str, Any]) -> str:
    display = score.get("display")
    if not isinstance(display, (int, float)):
        return "—"
    return f"{display:.1f}"


# Subscore/movement keys → leadership-facing labels (match the portal's
# since_last_audit vocabulary so the deck never renames a metric mid-report).
_DIMENSION_LABEL = {
    "visibility": "AI visibility",
    "attribution": "First-party citation",
    "category_visibility": "Category visibility",
    "identity": "Product identity",
    "content_richness": "Content depth",
    "routability": "Routability",
}


def _dim_label(key: Any, fallback: Any = None) -> str:
    k = str(key or "").strip().lower()
    if k in _DIMENSION_LABEL:
        return _DIMENSION_LABEL[k]
    return str(fallback or key or "").strip().replace("_", " ").title()


def _subscore_line(score: Mapping[str, Any]) -> Optional[str]:
    """"AI visibility 0.6 · First-party citation 4.6" — the composite score's
    breakdown so leadership sees WHICH dimension drags, not a bare number."""
    parts: List[str] = []
    for sub in _as_list(score.get("subscores")):
        sub = _as_dict(sub)
        disp = sub.get("display")
        if isinstance(disp, (int, float)):
            parts.append(f"{_dim_label(sub.get('key'))} {disp:.1f}")
    return "   ·   ".join(parts) if parts else None


def _material_movements(since: Mapping[str, Any]) -> List[str]:
    """Human trend lines from since_last_audit's MATERIAL movements only —
    "First-party citation 39 → 46 ▲". The old cover read score.delta (a
    different, usually-null field), so a re-audit deck showed no progress even
    when the report measured real movement."""
    out: List[str] = []
    _arrow = {"improved": "▲", "regressed": "▼"}
    for m in _as_list(since.get("movements")):
        m = _as_dict(m)
        if not m.get("is_material"):
            continue
        frm, to = m.get("from"), m.get("to")
        label = _dim_label(m.get("signal"), m.get("label"))
        arrow = _arrow.get(str(m.get("direction") or ""), "→")
        if isinstance(frm, (int, float)) and isinstance(to, (int, float)):
            out.append(f"{label} {frm:g} → {to:g} {arrow}")
        elif m.get("direction") == "changed":
            out.append(f"{label} changed")
    return out


def build_report_deck(
    summary: Mapping[str, Any],
    *,
    executive_bullets: Optional[List[str]] = None,
    preview_only: bool = False,
) -> Optional[bytes]:
    """Render the summary contract into PPTX bytes.

    preview_only=True (free tier) → a single watermarked cover slide: enough
    to share, honest about being a preview, and the upgrade hook. Returns None
    when python-pptx isn't installed (route surfaces a clear 503).
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Emu, Inches, Pt
    except ImportError:
        logger.warning("python-pptx not installed; deck export unavailable")
        return None

    summary = _as_dict(summary)
    score = _as_dict(summary.get("score"))
    verdict = _as_dict(summary.get("verdict"))
    meta = _as_dict(summary.get("meta"))
    subject = _as_dict(summary.get("subject"))
    brand = subject.get("merchant_name") or "Your brand"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def rgb(t):  # (r,g,b) -> RGBColor
        return RGBColor(*t)

    def add_slide(bg=_WHITE):
        slide = prs.slides.add_slide(blank)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb(bg)
        return slide

    def text_box(
        slide,
        x,
        y,
        w,
        h,
        lines,
        *,
        align=PP_ALIGN.LEFT,
    ):
        """lines: list of (text, size_pt, bold, color, space_after_pt)."""
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
        first = True
        for text, size, bold, color, space_after in lines:
            para = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            para.alignment = align
            para.space_after = Pt(space_after)
            run = para.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.name = "Calibri"
            run.font.color.rgb = rgb(color)
        return box

    band = str(score.get("band") or "")
    band_label, _band_color = _BAND_DECK_LABEL.get(band, ("", _MUTED))
    scale_max = score.get("scale_max") or 10
    generated_at = str(summary.get("generated_at") or "")[:10]

    # ── Cover: navy, giant score callout, verdict headline. ────────────────
    cover = add_slide(_NAVY)
    text_box(
        cover,
        0.9,
        0.7,
        11.5,
        0.6,
        [(f"{brand} — AI-readiness report", 20, True, _ICE, 0)],
    )
    text_box(
        cover,
        0.9,
        1.9,
        11.5,
        2.2,
        [
            (
                f"{_fmt_display(score)} / {scale_max}"
                + (f"   ·   {band_label}" if band_label else ""),
                66,
                True,
                _WHITE,
                6,
            )
        ],
    )
    # Subscore breakdown under the composite: the score is a weakest-link
    # number, so leadership needs to see which dimension it comes from.
    subscore_line = _subscore_line(score)
    if subscore_line:
        text_box(cover, 0.9, 3.25, 11.5, 0.5, [(subscore_line, 17, False, _ICE, 0)])
    cover_lines = []
    if verdict.get("headline"):
        cover_lines.append((str(verdict["headline"]), 20, False, _ICE, 10))
    # Real trend: material movements from since_last_audit (the old cover read
    # score.delta, a different, usually-null field, so re-audits showed no
    # progress even when the report measured it).
    trend = _material_movements(_as_dict(summary.get("since_last_audit")))
    if trend:
        cover_lines.append(("Since your last audit:  " + "    ".join(trend[:3]), 14, False, _ICE, 0))
    if cover_lines:
        text_box(cover, 0.9, 4.0, 11.5, 1.9, cover_lines)
    footer = f"Pivota AI-readiness audit · data as of {generated_at}" if generated_at else "Pivota AI-readiness audit"
    text_box(cover, 0.9, 6.7, 11.5, 0.4, [(footer, 11, False, _ICE, 0)])

    if preview_only:
        # Watermark + upgrade hook, then stop at one slide.
        text_box(
            cover,
            0.9,
            5.6,
            11.5,
            0.9,
            [
                (
                    "PREVIEW — upgrade your Pivota plan to export the full deck "
                    "(findings, competitive snapshot, and the action plan).",
                    16,
                    True,
                    _WHITE,
                    0,
                )
            ],
        )
        import io

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def bar(slide, x, y, w, h, frac, color):
        frac = max(0.0, min(1.0, float(frac)))
        # Track (full-width, faint) + fill (scaled). Zero-width fills crash
        # python-pptx, so a floor keeps a sliver visible.
        track = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        track.fill.solid(); track.fill.fore_color.rgb = rgb(_ICE); track.line.fill.background()
        fill_w = max(0.04, w * frac)
        f = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(fill_w), Inches(h))
        f.fill.solid(); f.fill.fore_color.rgb = rgb(color); f.line.fill.background()

    # ── Where you rank: Share of Voice — the competitive marquee. ──────────
    sov = _as_dict(summary.get("share_of_voice"))
    sov_brand = _as_dict(sov.get("brand"))
    sov_comps = [_as_dict(c) for c in _as_list(sov.get("competitors"))]
    if sov.get("available") and isinstance(sov_brand.get("pct"), (int, float)):
        slide = add_slide()
        text_box(slide, 0.9, 0.6, 11.5, 0.8, [("Where you rank in AI answers", 34, True, _NAVY, 0)])
        probed = sov.get("prompts_probed")
        sub = (
            f"Share of voice across {probed} category questions we tested — "
            "how often each brand is named when shoppers ask AI, not searching by name."
            if isinstance(probed, int)
            else "Share of voice — how often each brand is named when shoppers ask AI."
        )
        text_box(slide, 0.9, 1.35, 11.5, 0.5, [(sub, 13, False, _MUTED, 0)])
        # One ranked list including the brand, brand row highlighted. Sorted by
        # pct desc; brand pct = prompts_cited fraction.
        rows = [
            {"name": str(sov_brand.get("name") or brand), "pct": float(sov_brand.get("pct") or 0), "you": True}
        ]
        for c in sov_comps:
            if isinstance(c.get("pct"), (int, float)) and c.get("name"):
                rows.append({"name": str(c["name"]), "pct": float(c["pct"]), "you": False})
        rows.sort(key=lambda r: -r["pct"])
        you_rank = next((i + 1 for i, r in enumerate(rows) if r["you"]), None)
        if you_rank:
            text_box(
                slide, 0.9, 1.95, 11.5, 0.5,
                [(f"You rank #{you_rank} of {len(rows)} brands named.", 16, True, _NAVY, 0)],
            )
        y = 2.7
        for r in rows[:7]:
            color = _NAVY if r["you"] else _MUTED
            text_box(slide, 0.9, y - 0.02, 3.1, 0.4, [(r["name"] + ("  (you)" if r["you"] else ""), 13, r["you"], color, 0)])
            bar(slide, 4.1, y, 7.0, 0.28, r["pct"] / 100.0, _NAVY if r["you"] else _MUTED)
            text_box(slide, 11.3, y - 0.04, 1.1, 0.4, [(f"{r['pct']:.0f}%", 13, r["you"], color, 0)])
            y += 0.56

    # ── What moved since last audit: the outcome loop. ─────────────────────
    progress = _as_dict(summary.get("progress"))
    prog_sum = _as_dict(progress.get("summary"))
    prog_wins = [_as_dict(w) for w in _as_list(progress.get("wins"))]
    prog_prog = [_as_dict(p) for p in _as_list(progress.get("in_progress"))]
    if progress.get("available") and (
        any(prog_sum.get(k) for k in ("won", "progress", "no_change", "no_longer_grounded"))
    ):
        slide = add_slide()
        text_box(slide, 0.9, 0.6, 11.5, 0.8, [("What moved since last audit", 34, True, _NAVY, 0)])
        won, prog, nch = prog_sum.get("won", 0), prog_sum.get("progress", 0), prog_sum.get("no_change", 0)
        nlg = prog_sum.get("no_longer_grounded", 0)
        # A "what moved" slide must not hide regressions: surface lost
        # citations in both branches (review P3).
        dropped = f", {nlg} dropped off" if nlg else ""
        headline = (
            f"{won} host{'s' if won != 1 else ''} now recommend you, "
            f"{prog} started naming you, {nch} unchanged{dropped}."
            if won or prog else
            f"No new citations landed at your targets yet — {nch} unchanged"
            f"{dropped}. Working the moves below is what moves this."
        )
        color = _GREEN if (won or prog) else _MUTED
        text_box(slide, 0.9, 1.5, 11.5, 0.6, [(headline, 18, True, color, 0)])
        y = 2.4
        for row in (prog_wins + prog_prog)[:6]:
            wc = str(row.get("what_changed") or row.get("host") or "").strip()
            if wc:
                text_box(slide, 0.9, y, 11.5, 0.5, [(f"✓  {wc[:150]}", 14, False, _INK, 0)])
                y += 0.6
        note = str(progress.get("note") or "").strip()
        if note:
            text_box(slide, 0.9, 6.7, 11.5, 0.5, [(note[:180], 10, False, _MUTED, 0)])

    # ── What we found: findings + who AI cites instead. ────────────────────
    findings = [
        f for f in (_as_dict(x) for x in _as_list(summary.get("top_findings")))
        if f.get("title") or f.get("evidence_summary")
    ]
    snapshot = _as_dict(summary.get("competitive_snapshot"))
    slide = add_slide()
    text_box(slide, 0.9, 0.6, 11.5, 0.8, [("What we found", 34, True, _NAVY, 0)])
    y = 1.55
    # Lead with the score diagnosis: which measured dimension is the drag.
    weakest = _as_dict(score.get("weakest_dimension"))
    if isinstance(weakest.get("display"), (int, float)):
        text_box(
            slide, 0.9, y, 11.5, 0.5,
            [(
                f"Biggest drag on your score: {_dim_label(weakest.get('key'), weakest.get('label'))} "
                f"({weakest['display']:.1f}/{scale_max}).",
                16, True, _NAVY, 0,
            )],
        )
        y += 0.65
    for f in findings[:3]:
        sev_color = _SEVERITY_COLOR.get(str(f.get("severity") or ""), _MUTED)
        lines = []
        if f.get("title"):
            lines.append((str(f["title"]), 20, True, sev_color, 3))
        if f.get("evidence_summary"):
            lines.append((str(f["evidence_summary"]), 15, False, _INK, 0))
        text_box(slide, 0.9, y, 11.5, 1.3, lines)
        y += 1.45
    if snapshot.get("available"):
        snap_lines = [("Who AI cites instead", 16, True, _NAVY, 3)]
        hosts = [h for h in _as_list(snapshot.get("top_cited_hosts")) if h]
        comps = [c for c in _as_list(snapshot.get("competitors_named")) if c]
        if hosts:
            snap_lines.append(
                ("Sources: " + ", ".join(str(h) for h in hosts[:6]), 14, False, _INK, 2)
            )
        if comps:
            snap_lines.append(
                ("Competitors named: " + ", ".join(str(c) for c in comps[:6]), 14, False, _INK, 0)
            )
        if len(snap_lines) > 1:
            text_box(slide, 0.9, y, 11.5, 1.4, snap_lines)

    # ── Executive summary (LLM, optional). ──────────────────────────────────
    if executive_bullets:
        slide = add_slide()
        text_box(slide, 0.9, 0.6, 11.5, 0.8, [("Executive summary", 34, True, _NAVY, 0)])
        text_box(
            slide,
            0.9,
            1.8,
            11.5,
            4.6,
            [(f"•  {b}", 19, False, _INK, 14) for b in executive_bullets[:4]],
        )
        text_box(
            slide,
            0.9,
            6.8,
            11.5,
            0.4,
            [
                (
                    "Written from this audit's measured findings only.",
                    11,
                    False,
                    _MUTED,
                    0,
                )
            ],
        )

    # ── One slide per top action. ───────────────────────────────────────────
    all_actions = [
        a for a in (_as_dict(x) for x in _as_list(summary.get("top_actions")))
        if a.get("headline")
    ]
    # Leadership curation: lead with strategic PRIMARY moves; a re-test / QA
    # secondary ("Re-test failed SKU prompt: …") is a portal to-do, not a
    # boardroom slide. Fall back to all actions if curation would empty the
    # plan (never ship a deck with no moves).
    actions = [a for a in all_actions if a.get("action_source") != "secondary"] or all_actions
    for i, action in enumerate(actions[:3], start=1):
        slide = add_slide()
        text_box(
            slide,
            0.9,
            0.6,
            11.5,
            1.2,
            [(f"Move {i} — {action['headline']}", 30, True, _NAVY, 0)],
        )
        body = []
        if action.get("why_this_first"):
            body.append(("Why this first", 15, True, _MUTED, 2))
            body.append((str(action["why_this_first"]), 16, False, _INK, 10))
        if action.get("first_move"):
            body.append(("First move", 15, True, _MUTED, 2))
            body.append((str(action["first_move"]), 16, False, _INK, 0))
        if body:
            text_box(slide, 0.9, 2.0, 11.5, 3.0, body)
        prompts = [
            p for p in (_as_dict(x) for x in _as_list(action.get("supporting_prompts")))
            if p.get("query")
        ]
        if prompts and (action.get("supporting_prompts_basis") or "none") != "none":
            ev = [("What the AI answers showed", 13, True, _MUTED, 3)]
            for p in prompts[:3]:
                bits = [f"“{p['query']}”"]
                if p.get("provider"):
                    bits.append(str(p["provider"]))
                if p.get("reason"):
                    bits.append(str(p["reason"]))
                ev.append(("  ·  ".join(bits), 12, False, _MUTED, 2))
            text_box(slide, 0.9, 5.3, 11.5, 1.7, ev)

    # ── Where to earn citations: the off-platform channel moves. ───────────
    cited_moves = [
        _as_dict(m) for m in _as_list(summary.get("get_cited_moves"))
        if _as_dict(m).get("host")
    ]
    if cited_moves:
        slide = add_slide()
        text_box(slide, 0.9, 0.6, 11.5, 0.8, [("Where to earn citations", 34, True, _NAVY, 0)])
        text_box(
            slide, 0.9, 1.4, 11.5, 0.5,
            [("AI recommends independent sources, not your own page. These are the ones it "
              "grounds on for your category — earn a mention on each.", 13, False, _MUTED, 0)],
        )
        y = 2.15
        for m in cited_moves[:5]:
            host = str(m.get("host") or "")
            verb = "Build on" if m.get("already_endorses_you") else "Pitch"
            qs = [str(q) for q in _as_list(m.get("for_questions")) if q]
            head = f"{verb} {host}"
            if qs:
                head += f"  —  for “{qs[0]}”" + (f" +{len(qs)-1}" if len(qs) > 1 else "")
            text_box(slide, 0.9, y, 11.5, 0.4, [(head, 16, True, _INK, 2)])
            why = str(m.get("first_move") or m.get("why") or "").strip()
            if why:
                text_box(slide, 1.2, y + 0.4, 11.2, 0.5, [(why[:150], 12, False, _MUTED, 0)])
            y += 1.0
            if y > 6.6:
                break

    # ── Lanes you can win: the specific winnable demand. ───────────────────
    lanes = [
        _as_dict(l) for l in _as_list(summary.get("winnable_lanes"))
        if _as_dict(l).get("query")
    ]
    if lanes:
        slide = add_slide()
        text_box(slide, 0.9, 0.6, 11.5, 0.8, [("Lanes you can win", 34, True, _NAVY, 0)])
        text_box(
            slide, 0.9, 1.4, 11.5, 0.5,
            [("Specific shopper questions you're positioned to own — not the big-budget head "
              "terms. Win these first; each is tracked week over week.", 13, False, _MUTED, 0)],
        )
        y = 2.15
        _path_label = {"own_content": "Win with your own page", "publisher": "Get cited on a publisher"}
        for l in lanes[:5]:
            q = str(l.get("query") or "")
            text_box(slide, 0.9, y, 11.5, 0.4, [(f"“{q}”", 16, True, _INK, 2)])
            path = _path_label.get(str(l.get("win_path") or ""), "")
            hosts = [str(h) for h in _as_list(l.get("target_hosts")) if h]
            tail = path + (("  ·  via " + ", ".join(hosts[:3])) if hosts else "")
            if tail:
                text_box(slide, 1.2, y + 0.4, 11.2, 0.4, [(tail, 12, False, _MUTED, 0)])
            y += 0.92
            if y > 6.6:
                break

    # ── Per-product scorecard (multi-SKU only). ─────────────────────────────
    sku_rows = [
        s for s in (_as_dict(x) for x in _as_list(summary.get("sku_summaries")))
        if s.get("sku_title")
    ]
    if len(sku_rows) > 1:
        slide = add_slide()
        text_box(slide, 0.9, 0.6, 11.5, 0.8, [("Product-by-product", 34, True, _NAVY, 0)])
        text_box(
            slide, 0.9, 1.4, 11.5, 0.4,
            [("Each product's weakest-link score and status.", 13, False, _MUTED, 0)],
        )
        y = 2.1
        for s in sku_rows[:8]:
            sscore = _as_dict(s.get("score"))
            disp = _fmt_display(sscore)
            status = str(_as_dict(s.get("band_display")).get("label") or "")
            title = str(s.get("sku_title") or "")
            title = title[:58] + "…" if len(title) > 58 else title
            text_box(slide, 0.9, y, 7.3, 0.4, [(title, 14, True, _INK, 0)])
            text_box(slide, 8.4, y, 1.3, 0.4, [(f"{disp}/{scale_max}", 14, True, _NAVY, 0)])
            text_box(slide, 9.9, y, 2.5, 0.4, [(status, 12, False, _MUTED, 0)])
            y += 0.52
            if y > 6.6:
                break

    # ── Methodology / honesty close (navy bookend). ─────────────────────────
    slide = add_slide(_NAVY)
    text_box(slide, 0.9, 0.6, 11.5, 0.8, [("How this was measured", 34, True, _WHITE, 0)])
    limits = [str(l) for l in _as_list(meta.get("honest_limits")) if l]
    lines = [(f"•  {l}", 14, False, _ICE, 8) for l in limits[:5]]
    facts = []
    providers = [str(p) for p in _as_list(meta.get("providers")) if p]
    if providers:
        facts.append("Probed on " + ", ".join(providers))
    if isinstance(meta.get("products_audited"), int):
        facts.append(f"{meta['products_audited']} product(s) audited")
    if generated_at:
        facts.append(f"data as of {generated_at}")
    if facts:
        lines.append((" · ".join(facts), 13, False, _ICE, 0))
    if lines:
        text_box(slide, 0.9, 1.9, 11.5, 4.6, lines)
    text_box(
        slide,
        0.9,
        6.7,
        11.5,
        0.4,
        [("Generated by Pivota · pivota.cc", 12, False, _ICE, 0)],
    )

    import io

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
